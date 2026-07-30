#!/usr/bin/env python3
"""Build Safavieh CEO deck using Wayfair April 2026 template."""

from __future__ import annotations

import os
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

OUT_DIR = ROOT / "output" / "safavieh"
CHARTS_DIR = ROOT / "docs" / "small_parcel" / "safavieh_charts"
TEMPLATE = ROOT / "docs" / "small_parcel" / "Wayfair_Templates_April_2026.pptx"
PPTX_PATH = OUT_DIR / "Safavieh_CEO_June_MSBD.pptx"
CREDS_PATH = ROOT / ".gcp" / "credentials.json"

# Wayfair layout indices (after opening template)
LAYOUT_TITLE = 0       # TITLE
LAYOUT_SECTION = 1     # SECTION_HEADER
LAYOUT_BODY = 2        # TITLE_AND_BODY
LAYOUT_BLANK = 12      # BLANK


def delete_slide(prs: Presentation, index: int) -> None:
    slides = list(prs.slides._sldIdLst)
    slide_id = slides[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def trim_template(prs: Presentation) -> None:
    """Remove instruction and design-library slides from Wayfair template."""
    for _ in range(min(7, len(prs.slides))):
        delete_slide(prs, 0)
    while len(prs.slides) > 0:
        delete_slide(prs, 0)


def _set_title_and_body(slide, title: str, body: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title
    # body placeholder
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx != 0 and ph.has_text_frame:
            ph.text_frame.clear()
            for i, line in enumerate(body.split("\n")):
                p = ph.text_frame.paragraphs[0] if i == 0 else ph.text_frame.add_paragraph()
                p.text = line
                p.level = 0
                p.font.size = Pt(14)
            break


def _section_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    if slide.shapes.title:
        slide.shapes.title.text = title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            ph.text = subtitle
            break


def _body_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BODY])
    body = "\n".join(f"• {b}" for b in bullets)
    _set_title_and_body(slide, title, body)


def _chart_slide(prs: Presentation, image_path: Path, title: str) -> None:
    if not image_path.exists():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    # title textbox
    left, top, width, height = Inches(0.4), Inches(0.15), Inches(9), Inches(0.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.text = title
    box.text_frame.paragraphs[0].font.size = Pt(20)
    box.text_frame.paragraphs[0].font.bold = True
    slide.shapes.add_picture(str(image_path), Inches(0.35), Inches(0.65), width=Inches(9.3))


def build_presentation() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Wayfair template not found: {TEMPLATE}")

    shutil.copy(TEMPLATE, PPTX_PATH)
    prs = Presentation(str(PPTX_PATH))
    trim_template(prs)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    if slide.shapes.title:
        slide.shapes.title.text = "Safavieh CEO In-Office"
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            ph.text = "June MSBD performance · Badging simulation · Same-day induction\nWayfair Small Parcel · July 2026"
            break

    _section_slide(prs, "Executive summary", "June 2026 MSBD · Rugs dropship · 73k ops")

    _body_bullets(
        prs,
        "Key takeaways",
        [
            "90% IFR in June — strong month; 68.3% same-day induction before 2pm (Mon–Fri)",
            "Fri/Sat −1 o2d (Sun MSBD): +9.0 pp at 3-day · 6,595 newly badged orders vs current stated",
            "Full policy stack: 6.8% 1-day · 27.1% 2-day · 60.2% 3-day (includes cutoff + Fri/Sat −1)",
            "Cutoff extension: weekdays only (Mon–Fri), after cutoff before 2pm — 5,557 orders",
            "Ops check: 68.3% same-day induct before 2pm · ~90% IFR — can they deliver faster badges?",
        ],
    )

    _chart_slide(prs, CHARTS_DIR / "04_badging_tiers_current_vs_sim.png", "Badging by speed tier — June vs full simulation")
    _chart_slide(prs, CHARTS_DIR / "08_weekend_incremental_by_tier.png", "Sunday MSBD — Fri/Sat −1 o2d lift vs current stated")

    _body_bullets(
        prs,
        "Weekend shipping — Sunday MSBD promise",
        [
            "Fri/Sat placed → o2d_stated − 1 vs current stated (17,881 orders in June)",
            "Lift vs current: +1.8 pp 1-day · +3.7 pp 2-day · +9.0 pp 3-day · +1.9 pp fast",
            "6,595 additional 3-day badges · 1,388 additional fast badges",
            "~33% of Fri/Sat orders induct Sat/Sun in June (execution; separate from stated lift)",
        ],
    )

    _chart_slide(prs, CHARTS_DIR / "01_ifr_by_warehouse.png", "IFR by warehouse (June MSBD)")
    _chart_slide(prs, CHARTS_DIR / "02_before_2pm_same_day_induction.png", "Same-day induction before 2pm by warehouse")
    _chart_slide(prs, CHARTS_DIR / "03_ifr_vs_before_2pm_scatter.png", "IFR vs before-2pm induction")
    _chart_slide(prs, CHARTS_DIR / "05_badging_opportunity_uplift.png", "Total badging opportunity by tier")
    _chart_slide(prs, CHARTS_DIR / "06_3d_badge_by_warehouse.png", "3-day badge by warehouse")

    _section_slide(prs, "Discussion questions", "CEO conversation")

    _body_bullets(
        prs,
        "Questions for Safavieh",
        [
            "Prioritize 2pm + no cushion first — drives most 2- and 3-day badge lift",
            "Weekend stated promise: +9.0 pp at 3-day · 6,595 newly badged orders (Fri/Sat −1 vs current)",
            "Can NJ/PA sites improve before-2pm induction (42–50% vs 94% Carlisle)?",
            "FedEx pickup schedule vs proposed 2pm window at each DC?",
        ],
    )

    prs.save(PPTX_PATH)
    return PPTX_PATH


def upload_to_google_slides(pptx_path: Path) -> str | None:
    if not CREDS_PATH.exists():
        return None
    os.environ.setdefault("GOOGLE_APPLICATION_CREDENTIALS", str(CREDS_PATH))
    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload

    scopes = [
        "https://www.googleapis.com/auth/drive",
        "https://www.googleapis.com/auth/presentations",
    ]
    creds = service_account.Credentials.from_service_account_file(str(CREDS_PATH), scopes=scopes)
    drive = build("drive", "v3", credentials=creds)
    media = MediaFileUpload(
        str(pptx_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )
    meta = {"name": "Safavieh CEO — June MSBD Analysis", "mimeType": "application/vnd.google-apps.presentation"}
    file = drive.files().create(body=meta, media_body=media, fields="id,webViewLink").execute()
    drive.permissions().create(fileId=file["id"], body={"type": "anyone", "role": "reader"}).execute()
    return file.get("webViewLink") or f"https://docs.google.com/presentation/d/{file['id']}/edit"


def main() -> None:
    path = build_presentation()
    print(f"Created {path} (Wayfair template)")
    try:
        link = upload_to_google_slides(path)
        if link:
            print(f"Google Slides: {link}")
    except Exception as e:
        print(f"Google upload not available ({e})")
        print("Import PPTX to Google Slides: File → Import slides")


if __name__ == "__main__":
    main()
